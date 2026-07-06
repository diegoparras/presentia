"""Scene JSON -> native, editable PPTX. Reports native/approx/fallback coverage."""
import sys, json, collections
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

EMU = 9525  # 1 CSS px @96dpi
def px(v): return Emu(int(v * EMU))
def pt(v): return Pt(v * 0.75)
def rgb(h):
    h = (h or "#000000").lstrip("#"); return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

scene = json.load(open(sys.argv[1]))
out = sys.argv[2]
cov = collections.Counter()
detail = []

prs = Presentation(); prs.slide_width = px(scene["width"]); prs.slide_height = px(scene["height"])
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes

for b in scene["blocks"]:
    t = b["type"]
    if t == "rect":
        approx = b.get("gradient") or b.get("shadow")
        shp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if b.get("radius") else MSO_SHAPE.RECTANGLE,
                               px(b["x"]), px(b["y"]), px(b["w"]), px(b["h"]))
        shp.fill.solid(); shp.fill.fore_color.rgb = rgb(b.get("fill"))
        shp.line.fill.background()
        shp.shadow.inherit = False
        cov["approx" if approx else "native"] += 1
        detail.append(("rect", "APROX (gradiente/sombra→sólido)" if approx else "NATIVO"))
    elif t == "text":
        tb = shapes.add_textbox(px(b["x"]), px(b["y"]), px(max(b["w"],120)), px(max(b["h"],30)))
        tf = tb.text_frame; tf.word_wrap = True; tf.margin_left=0; tf.margin_top=0; tf.margin_right=0; tf.margin_bottom=0
        p = tf.paragraphs[0]; p.alignment = ALIGN.get(b.get("align"), PP_ALIGN.LEFT)
        r = p.add_run(); r.text = b["text"]
        r.font.size = pt(b["fontSize"]); r.font.bold = bool(b.get("bold")); r.font.color.rgb = rgb(b.get("color"))
        cov["native"] += 1; detail.append(("text", "NATIVO"))
    elif t == "chart":
        cd = CategoryChartData(); cd.categories = b["data"]["categories"]
        for s in b["data"]["series"]:
            cd.add_series(s["name"], s["values"])
        xl = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "line": XL_CHART_TYPE.LINE,
              "pie": XL_CHART_TYPE.PIE, "area": XL_CHART_TYPE.AREA}.get(b.get("chartType"), XL_CHART_TYPE.COLUMN_CLUSTERED)
        shapes.add_chart(xl, px(b["x"]), px(b["y"]), px(b["w"]), px(b["h"]), cd)
        cov["native"] += 1; detail.append(("chart", "NATIVO (editable)"))
    elif t == "table":
        hdr = b["table"]["headers"]; rows = b["table"]["rows"]
        gf = shapes.add_table(len(rows)+1, len(hdr), px(b["x"]), px(b["y"]), px(b["w"]), px(40*(len(rows)+1)))
        tbl = gf.table
        for c, htxt in enumerate(hdr):
            cell = tbl.cell(0, c); cell.text = htxt
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb("#7c51f8")
            rp = cell.text_frame.paragraphs[0].runs[0]; rp.font.bold=True; rp.font.color.rgb=rgb("#ffffff"); rp.font.size=pt(18)
        for ri, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(ri, c); cell.text = str(val)
                rp = cell.text_frame.paragraphs[0].runs[0]; rp.font.size=pt(16); rp.font.color.rgb=rgb("#374151")
        cov["native"] += 1; detail.append(("table", "NATIVO (editable)"))
    else:
        cov["fallback"] += 1; detail.append((t, "FALLBACK-IMAGEN"))

prs.save(out)
total = sum(cov.values())
print(f"\n=== COBERTURA ({out.split('/')[-1]}) — {total} bloques ===")
for k in ("native","approx","fallback"):
    if cov[k]: print(f"  {k:9}: {cov[k]:2}  ({cov[k]/total*100:.0f}%)")
for typ, verd in detail:
    print(f"    - {typ:6} {verd}")
