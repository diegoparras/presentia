"""Frozen scene IR -> PPTX (no headless browser at build time).

Reads the JSON produced by the freeze extractor (a list of {html, scene[, png]})
and builds a .pptx where text, rectangles and images are **native, editable**
PowerPoint objects (real textboxes / autoshapes / pictures), not a flat
screenshot of the whole slide. Charts and vector art are embedded as the crisp
PNGs captured during the single freeze pass.

This is the structural inverse of the current PPTX export (which rasterizes each
whole slide): here everything except the charts stays editable. Native chart
reconstruction (data -> real PPTX chart) is a per-template follow-up.
"""
from __future__ import annotations

import base64
import collections
import io
import json
import re
import sys
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

EMU_PER_PX = 9525  # 1 CSS px @96dpi
ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}


def px(v: float) -> Emu:
    return Emu(int(round(v * EMU_PER_PX)))


def pt(v: float) -> Pt:
    # CSS px -> pt (72/96).
    return Pt(v * 0.75)


def to_rgb(color: str | None) -> RGBColor:
    if not color:
        return RGBColor(0, 0, 0)
    m = re.match(r"rgba?\(([^)]+)\)", color)
    if m:
        parts = [p.strip() for p in m.group(1).split(",")]
        r, g, b = (int(float(parts[i])) for i in range(3))
        return RGBColor(r, g, b)
    h = color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) == 6:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return RGBColor(0, 0, 0)


def is_opaque(color: str | None) -> bool:
    if not color:
        return False
    m = re.match(r"rgba\(([^)]+)\)", color)
    if m:
        parts = [p.strip() for p in m.group(1).split(",")]
        return len(parts) < 4 or float(parts[3]) > 0.05
    return color not in ("transparent",)


def add_text(shapes, b: dict[str, Any]) -> None:
    tb = shapes.add_textbox(px(b["x"]), px(b["y"]), px(max(b["w"], 8)), px(max(b["h"], 8)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = ALIGN.get(b.get("align"), PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = b.get("text", "")
    run.font.size = pt(b.get("fontSize", 16))
    run.font.bold = bool(b.get("bold"))
    run.font.color.rgb = to_rgb(b.get("color"))
    fam = (b.get("fontFamily") or "").split(",")[0].strip().strip("'\"")
    if fam:
        run.font.name = fam


def add_rect(shapes, b: dict[str, Any]) -> None:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if b.get("radius") else MSO_SHAPE.RECTANGLE
    shp = shapes.add_shape(shape_type, px(b["x"]), px(b["y"]), px(b["w"]), px(b["h"]))
    if is_opaque(b.get("fill")):
        shp.fill.solid()
        shp.fill.fore_color.rgb = to_rgb(b["fill"])
    else:
        shp.fill.background()
    shp.line.fill.background()
    shp.shadow.inherit = False


def add_png(shapes, b: dict[str, Any]) -> bool:
    if not b.get("png"):
        return False
    data = base64.b64decode(b["png"])
    shapes.add_picture(io.BytesIO(data), px(b["x"]), px(b["y"]), px(b["w"]), px(b["h"]))
    return True


def add_image(shapes, b: dict[str, Any]) -> bool:
    src = b.get("src", "")
    if src.startswith("data:image"):
        header, _, payload = src.partition(",")
        try:
            shapes.add_picture(io.BytesIO(base64.b64decode(payload)), px(b["x"]), px(b["y"]), px(b["w"]), px(b["h"]))
            return True
        except Exception:
            return False
    # Remote/relative images: skip embedding here (no network at build time);
    # leave a placeholder rect so geometry is preserved.
    return False


def build_pptx(slides: list[dict[str, Any]], out_path: str) -> collections.Counter:
    if not slides:
        raise SystemExit("no slides")
    w = slides[0]["scene"]["width"]
    h = slides[0]["scene"]["height"]
    prs = Presentation()
    prs.slide_width = px(w)
    prs.slide_height = px(h)
    blank = prs.slide_layouts[6]
    cov: collections.Counter = collections.Counter()

    for slide in slides:
        s = prs.slides.add_slide(blank)
        shapes = s.shapes
        for b in slide["scene"]["blocks"]:
            t = b["type"]
            if t == "rect":
                add_rect(shapes, b)
                cov["rect-native"] += 1
            elif t == "text":
                add_text(shapes, b)
                cov["text-native"] += 1
            elif t == "svg":
                if add_png(shapes, b):
                    cov["chart-image"] += 1
                else:
                    cov["chart-missing"] += 1
            elif t == "image":
                if add_image(shapes, b):
                    cov["image-native"] += 1
                else:
                    cov["image-skipped"] += 1

    prs.save(out_path)
    return cov


def main() -> None:
    frozen_path = sys.argv[1]
    out_path = sys.argv[2]
    with open(frozen_path, encoding="utf-8") as fh:
        slides = json.load(fh)
    cov = build_pptx(slides, out_path)
    total = sum(cov.values())
    native = cov["rect-native"] + cov["text-native"] + cov["image-native"]
    print(f"pptx: {len(slides)} slides -> {out_path}")
    print(f"  bloques: {total}  nativos-editables: {native}  charts-imagen: {cov['chart-image']}")
    print(f"  detalle: {dict(cov)}")


if __name__ == "__main__":
    main()
